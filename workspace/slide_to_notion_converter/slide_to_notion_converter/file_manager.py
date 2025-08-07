import os
import json
from typing import Dict
from pptx import Presentation

class FileManager:
    """
    This class is responsible for reading and writing files.
    """
    def __init__(self):
        pass

    def read_file(self, file_path: str) -> Dict:
        """
        Read the PowerPoint file and return the slide data.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"The file {file_path} does not exist.")

        slide_data = {}
        presentation = Presentation(file_path)
        for i, slide in enumerate(presentation.slides):
            slide_data[i] = [shape.text for shape in slide.shapes if shape.has_text_frame]

        return slide_data

    def write_file(self, file_path: str, data: Dict):
        """
        Write the Notion data to a file.
        """
        with open(file_path, 'w') as f:
            json.dump(data, f)
